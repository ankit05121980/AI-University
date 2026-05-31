"""Render a 30+ slide executive presentation deck for a :class:`Book` using
python-pptx, including speaker notes, an executive summary, a technology
overview and architecture content.
"""
from __future__ import annotations

import textwrap

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from ..models import Book

ACCENT = RGBColor(0x43, 0x38, 0xCA)
DARK = RGBColor(0x0F, 0x17, 0x2A)
MUTED = RGBColor(0x47, 0x55, 0x69)


def _first_sentences(text: str, n: int = 2) -> str:
    parts = [p.strip() for p in text.replace("\n", " ").split(". ") if p.strip()]
    out = ". ".join(parts[:n])
    if out and not out.endswith("."):
        out += "."
    return out


def _bullets(text: str, n: int = 4) -> list[str]:
    parts = [p.strip() for p in text.replace("\n", " ").split(". ") if len(p.strip()) > 25]
    return [(p if p.endswith(".") else p + ".") for p in parts[:n]]


def render_pptx(book: Book, path: str) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def add_slide(title: str, bullets: list[str], notes: str = "", *,
                  subtitle: str | None = None):
        slide = prs.slides.add_slide(blank)
        # accent bar
        bar = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(13.333), Inches(0.18))
        bar.fill  # no-op to keep linter calm
        # title
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(30)
        run.font.bold = True
        run.font.color.rgb = ACCENT
        if subtitle:
            sp = tf.add_paragraph()
            sr = sp.add_run()
            sr.text = subtitle
            sr.font.size = Pt(15)
            sr.font.color.rgb = MUTED
        # body bullets
        if bullets:
            body = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(12.0),
                                            Inches(5.2))
            bf = body.text_frame
            bf.word_wrap = True
            for i, b in enumerate(bullets):
                para = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
                r = para.add_run()
                r.text = "\u2022 " + b
                r.font.size = Pt(18)
                r.font.color.rgb = DARK
                para.space_after = Pt(8)
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        return slide

    # Title slide
    s = prs.slides.add_slide(blank)
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(3))
    tf = box.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = book.title
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = ACCENT
    sp = tf.add_paragraph()
    sr = sp.add_run()
    sr.text = book.subtitle
    sr.font.size = Pt(18)
    sr.font.color.rgb = MUTED
    mp = tf.add_paragraph()
    mr = mp.add_run()
    mr.text = f"AI-University Press   |   {', '.join(book.authors)}   |   {book.published}"
    mr.font.size = Pt(13)
    mr.font.color.rgb = MUTED
    s.notes_slide.notes_text_frame.text = (
        f"Executive presentation for '{book.title}'. This deck summarises the book for "
        "leadership and stakeholders.")

    fm = book.front_matter
    add_slide("Executive Summary", _bullets(fm.executive_summary, 5),
              notes=fm.executive_summary)
    add_slide("Learning Objectives", fm.learning_objectives[:6],
              notes="Walk through what the audience will be able to do.")
    add_slide("Industry Context", _bullets(fm.industry_context, 4),
              notes=fm.industry_context)
    add_slide("Business Perspective", _bullets(fm.business_perspective, 4),
              notes=fm.business_perspective)
    add_slide("Technology Overview", _bullets(fm.technical_perspective, 5),
              notes=fm.technical_perspective)
    add_slide("Architecture Perspective", _bullets(fm.architecture_perspective, 4),
              notes=fm.architecture_perspective)

    # Chapter highlight slides (cap so we comfortably exceed 30 total)
    for ch in book.chapters:
        bullets = [_first_sentences(ch.summary, 2)]
        # add a couple of section headlines
        bullets += [s.heading for s in ch.sections[:3] if s.heading not in
                    ("Introduction",)][:3]
        add_slide(f"Chapter {ch.number}: {ch.title}", bullets[:5],
                  notes=ch.summary, subtitle=book.category)

    add_slide("Governance & Responsible AI", _bullets(fm.governance_perspective, 4),
              notes=fm.governance_perspective)
    add_slide("Security Perspective", _bullets(fm.security_perspective, 4),
              notes=fm.security_perspective)

    # Case study slide
    if book.back_matter.case_studies:
        cs = book.back_matter.case_studies[0]
        add_slide(cs["title"], _bullets(cs["body"], 4), notes=cs["body"])

    add_slide("Key Takeaways", [
        "Treat AI as a systems discipline: data, models and operations together.",
        "Measure everything; gate releases on evaluation, not intuition.",
        "Design security, governance and cost controls in from the start.",
        "Prefer the simplest design that meets the requirement.",
        "Operate continuously: monitor, learn and improve.",
    ], notes="Close with the durable principles from the book.")

    add_slide("Thank You", [
        f"{book.title}",
        "AI-University Press — Enterprise AI Knowledge Series",
        f"ISBN {book.isbn}  |  v{book.version}",
    ], notes="Q&A and next steps.")

    prs.save(path)
